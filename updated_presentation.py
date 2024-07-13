import csv
import os
from pptx import Presentation
from pptx.util import Inches

def load_csv(csv_file):
    data = {}
    with open(csv_file, mode='r') as file:
        reader = csv.DictReader(file) 
        for row in reader:
            data[row['Key']] = row['Data']
    print("CSV Data:", data)  
    return data

def update_pptx(pptx_file, data, output_file):
    prs = Presentation(pptx_file)
    for slide in prs.slides:
        for shape in slide.shapes: 
            if shape.has_text_frame: 
                for paragraph in shape.text_frame.paragraphs: 
                    for run in paragraph.runs: 
                        key = run.text.strip() 
                        if key in data: 
                            run.text = data[key] 
                            print(f"Updated text for {key}: {data[key]}") 
            elif shape.shape_type == 13: 
                key = shape.name.strip()
                if key in data: #if the key is in the dictionary
                    image_path = data[key] #check the image form the dictionary
                    if os.path.exists(image_path):
                        shape._element.getparent().remove(shape._element)  
                        slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
                        print(f"Updated image for {key}: {image_path}") 
                    else:
                        print(f"Image path does not exist: {image_path}")
    prs.save(output_file)

csv_file = 'data.csv'
pptx_file = 'source_presentation.pptx'
output_file = r'/home/dev/Desktop/updated_presentation.pptx'

if not os.path.exists(csv_file):
    raise FileNotFoundError(f"The file {csv_file} does not exist.")
if not os.path.exists(pptx_file):
    raise FileNotFoundError(f"The file {pptx_file} does not exist.")

if os.path.exists(output_file):
    os.remove(output_file)

csv_data = load_csv(csv_file)
update_pptx(pptx_file, csv_data, output_file)
